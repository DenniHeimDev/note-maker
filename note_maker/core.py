from __future__ import annotations

import os
import shutil
from dataclasses import dataclass
from pathlib import Path
from typing import Optional

import fitz  # PyMuPDF
from dotenv import load_dotenv
from openai import OpenAI
from pptx import Presentation

from config_helpers import USER_CONFIG_DIR

load_dotenv()

AVAILABLE_MODELS = [
    "gpt-5",
    "gpt-5.1",
    "gpt-4.1",
    "gpt-4.1-mini",
    "gpt-5-mini",
]
DEFAULT_MODEL = "gpt-5.1"

LANGUAGE_OPTIONS = {
    "nynorsk": {
        "label": "Nynorsk",
        "system_prompt": (
            "Du er ein fagleg dyktig skribent som skriv klart og presist på nynorsk.\n"
            "Du får tekst frå ei fagleg presentasjon (PowerPoint eller PDF) og skal lage\n"
            "eit strukturert notat på nynorsk. Behald fagterminologi, ikkje oversett direkte til norsk vist det står på engelsk, bruk overskrifter\n"
            "og underoverskrifter der det passar, og skriv i ein stil som eignar seg\n"
            "som førebuing til undervisning eller eksamen. Ikkje omslutt notatet i ```markdown```-blokker."
        ),
        "user_template": (
            "Her er innhaldet frå presentasjonen. Lag eit strukturert notat på nynorsk\n"
            "som oppsummerer og forklarer innhaldet. Du skal ikkje referere til \"slides\"\n"
            "eller \"bilete\", berre skrive eit samanhengande notat i markdown-format utan å bruke ```markdown``` eller andre kodeblokker.\n\n"
            "=== START AV INPUT ===\n"
            "{tekst_her}\n"
            "=== SLUTT AV INPUT ==="
        ),
        "note_suffix": "notat_nynorsk",
    },
    "bokmal": {
        "label": "Bokmål",
        "system_prompt": (
            "Du er en faglig dyktig skribent som skriver klart og presist på bokmål.\n"
            "Du får tekst fra en faglig presentasjon (PowerPoint eller PDF) og skal lage\n"
            "et strukturert notat på bokmål. Behold fagterminologi, ikke oversett direkte fra engelsk\n"
            "dersom det ikke gir mening, og bruk overskrifter og underoverskrifter der det passer. Ikke bruk ```markdown```-blokker rundt notatet."
        ),
        "user_template": (
            "Her er innholdet fra presentasjonen. Lag et strukturert notat på bokmål\n"
            "som oppsummerer og forklarer innholdet. Du skal ikke referere til \"slides\"\n"
            "eller \"bilder\", men skrive et sammenhengende notat i markdown uten å omslutte teksten med ```markdown```.\n\n"
            "=== START AV INPUT ===\n"
            "{tekst_her}\n"
            "=== SLUTT AV INPUT ==="
        ),
        "note_suffix": "notat_bokmal",
    },
    "english": {
        "label": "English",
        "system_prompt": (
            "You are an expert technical writer who produces clear, structured notes in English.\n"
            "You receive text extracted from a presentation (PowerPoint or PDF) and must create\n"
            "a study note. Keep domain terminology, avoid literal translations that harm meaning,\n"
            "and use headings and subheadings where appropriate to prepare the reader for teaching or exams. Never wrap the output in ```markdown``` code fences."
        ),
        "user_template": (
            "Here is the content from the presentation. Produce a structured note in English\n"
            "that summarizes and explains the material. Do not mention \"slides\" or \"images\";\n"
            "write a continuous markdown document instead, but do not surround the note with ```markdown``` fences.\n\n"
            "=== START OF INPUT ===\n"
            "{tekst_her}\n"
            "=== END OF INPUT ==="
        ),
        "note_suffix": "note_english",
    },
}
DEFAULT_LANGUAGE = "nynorsk"
LANGUAGE_LABEL_TO_KEY = {data["label"]: key for key, data in LANGUAGE_OPTIONS.items()}
LANGUAGE_LABELS = [data["label"] for data in LANGUAGE_OPTIONS.values()]

_CLIENT: Optional[OpenAI] = None
_API_KEY = None


def _ensure_api_key() -> Optional[str]:
    """
    Returns the OpenAI API key from the environment, falling back to plain-text .env files.
    Supports both standard KEY=VALUE files and files that only contain the key value.
    """
    existing = os.environ.get("OPENAI_API_KEY")
    if existing:
        return existing

    candidate_paths = [
        Path(".env"),
        Path(__file__).with_name(".env"),
        Path(__file__).resolve().parent.parent / ".env",
        USER_CONFIG_DIR / "config.env",
    ]

    for path in candidate_paths:
        if not path.exists():
            continue
        with path.open(encoding="utf-8") as handle:
            for line in handle:
                stripped = line.strip()
                if not stripped or stripped.startswith("#"):
                    continue
                if "=" in stripped:
                    _, value = stripped.split("=", 1)
                    candidate = value.strip().strip("'\"")
                else:
                    candidate = stripped
                if candidate:
                    os.environ["OPENAI_API_KEY"] = candidate
                    return candidate
    return None


def _get_client() -> OpenAI:
    global _CLIENT, _API_KEY
    api_key = _ensure_api_key()
    if not api_key:
        raise RuntimeError("Miljøvariabelen OPENAI_API_KEY er ikkje sett.")
    if _CLIENT is None or api_key != _API_KEY:
        _CLIENT = OpenAI(api_key=api_key)
        _API_KEY = api_key
    return _CLIENT


def extract_text_from_pptx(file_path: str, include_notes: bool = False) -> str:
    """Extract text content from a PowerPoint presentation.

    We preserve a bit of structure (section separators, title hints, bullet levels)
    to give the model higher-quality input.
    """

    def _iter_text_frame_lines(text_frame) -> list[str]:
        lines: list[str] = []
        for paragraph in getattr(text_frame, "paragraphs", []) or []:
            text = getattr(paragraph, "text", "")
            text = text.strip() if text else ""
            if not text:
                continue
            level = int(getattr(paragraph, "level", 0) or 0)
            indent = "  " * max(level, 0)
            # Represent paragraphs as bullets to keep hierarchy.
            lines.append(f"{indent}- {text}")
        return lines

    presentation = Presentation(file_path)
    sections: list[str] = []

    for idx, slide in enumerate(presentation.slides, start=1):
        slide_parts: list[str] = []

        # Try to extract a title (if the presentation uses a title placeholder).
        title_text = ""
        title_shape = None
        try:
            title_shape = slide.shapes.title
        except Exception:
            title_shape = None

        if title_shape is not None and getattr(title_shape, "has_text_frame", False):
            title_text = title_shape.text_frame.text.strip() if title_shape.text_frame.text else ""

        # Collect slide content.
        content_lines: list[str] = []
        for shape in slide.shapes:
            if title_shape is not None and shape is title_shape:
                continue

            if getattr(shape, "has_text_frame", False):
                content_lines.extend(_iter_text_frame_lines(shape.text_frame))
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        content_lines.append(" | ".join(cells))

        notes_text = ""
        if include_notes:
            try:
                if slide.has_notes_slide and slide.notes_slide and slide.notes_slide.notes_text_frame:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
            except Exception:
                notes_text = ""

        if not title_text and not content_lines and not notes_text:
            continue

        slide_parts.append(f"=== SECTION {idx} ===")
        if title_text:
            slide_parts.append(f"TITLE: {title_text}")
        if content_lines:
            slide_parts.extend(content_lines)
        if include_notes and notes_text:
            slide_parts.append("NOTES:")
            slide_parts.append(notes_text)

        sections.append("\n".join(slide_parts))

    return "\n\n".join(sections)


def _extract_tables_from_pdf_page(page) -> list[str]:
    """Best-effort table extraction from a PDF page.

    Uses PyMuPDF's `page.find_tables()` when available. Output is intentionally
    simple (pipe-separated rows) to keep it robust across PDFs.
    """

    if not hasattr(page, "find_tables"):
        return []

    try:
        finder = page.find_tables()
    except Exception:
        return []

    tables = getattr(finder, "tables", None) or []
    lines: list[str] = []

    for table_index, table in enumerate(tables, start=1):
        extract = getattr(table, "extract", None)
        if not callable(extract):
            continue
        try:
            rows = extract()
        except Exception:
            continue
        if not rows:
            continue

        lines.append(f"TABLE {table_index}:")
        for row in rows:
            if not row:
                continue
            cells = [str(cell).strip() for cell in row if cell is not None and str(cell).strip()]
            if cells:
                lines.append(" | ".join(cells))

    return lines


def extract_text_from_pdf(file_path: str) -> str:
    """Extract text content from a PDF file.

    Adds page markers when the document has multiple pages to preserve structure.
    Also attempts best-effort table extraction.
    """

    doc = fitz.open(file_path)
    pages: list[str] = []
    try:
        multi_page = doc.page_count > 1
        for idx, page in enumerate(doc, start=1):
            page_text = page.get_text("text").strip()
            table_lines = _extract_tables_from_pdf_page(page)

            if not page_text and not table_lines:
                continue

            parts: list[str] = []
            if multi_page:
                parts.append(f"=== PAGE {idx} ===")
            if page_text:
                parts.append(page_text)
            if table_lines:
                parts.append("TABLES:")
                parts.extend(table_lines)

            pages.append("\n".join(parts))
    finally:
        doc.close()

    return "\n\n".join(pages)


def extract_text(file_path: str, include_pptx_notes: bool = False) -> str:
    """Dispatch text extraction based on file extension."""

    suffix = Path(file_path).suffix.lower()
    if suffix == ".pptx":
        return extract_text_from_pptx(file_path, include_notes=include_pptx_notes)
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    raise ValueError("Ukjend filtype. Vel ei .pptx- eller .pdf-fil.")


def _language_settings(language_key: str) -> dict:
    """Retrieve language-specific settings."""
    settings = LANGUAGE_OPTIONS.get(language_key)
    if not settings:
        raise ValueError(f"Ugyldig språkval: {language_key}.")
    return settings


def generate_note_from_text(text: str, model_name: str, language_key: str) -> str:
    """Generate a study note from the provided text using OpenAI."""
    if not text.strip():
        raise ValueError("Fann ikkje tekst i den valde fila.")
    if model_name not in AVAILABLE_MODELS:
        raise ValueError(f"Ugyldig modell: {model_name}.")
    client = _get_client()
    settings = _language_settings(language_key)
    user_prompt = settings["user_template"].format(tekst_her=text.strip())
    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "system", "content": settings["system_prompt"]},
            {"role": "user", "content": user_prompt},
        ],
    )
    note = response.choices[0].message.content
    if not note:
        raise RuntimeError("Modellen returnerte ikkje noko innhald.")
    return note.strip()


def save_note_text(note_text: str, output_dir: Path | str, source_name: str, note_suffix: str) -> Path:
    """Save the generated note to a markdown file."""
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    file_stem = Path(source_name).stem
    note_file = output_path / f"{file_stem}_{note_suffix}.md"
    with note_file.open("w", encoding="utf-8") as handle:
        handle.write(note_text)
    return note_file


def copy_source_file(
    source_file: Path | str,
    output_dir: Path | str,
    desired_name: Optional[str] = None,
) -> Path:
    """Copy the source file to the output directory, handling duplicates."""
    destination_dir = Path(output_dir)
    destination_dir.mkdir(parents=True, exist_ok=True)

    name_source = Path(desired_name) if desired_name else Path(source_file)
    stem = name_source.stem
    suffix = name_source.suffix
    destination = destination_dir / f"{stem}{suffix}"

    counter = 1
    while destination.exists():
        destination = destination_dir / f"{stem}_kopi_{counter}{suffix}"
        counter += 1

    shutil.copy2(source_file, destination)
    return destination


@dataclass
class GenerationResult:
    note_path: Path
    note_text: str
    copied_path: Optional[Path]


def generate_note_from_file(
    source_path: Path,
    original_filename: str,
    output_dir: Path,
    model_name: str,
    language_key: str,
    copy_requested: bool = False,
    copy_dir: Optional[Path] = None,
    include_pptx_notes: bool = False,
) -> GenerationResult:
    """Orchestrate the note generation process from a file."""
    extracted_text = extract_text(str(source_path), include_pptx_notes=include_pptx_notes)
    note_text = generate_note_from_text(extracted_text, model_name, language_key)
    language_settings = _language_settings(language_key)
    note_path = save_note_text(
        note_text,
        output_dir,
        original_filename,
        language_settings["note_suffix"],
    )
    copied_path = None
    if copy_requested:
        target_dir = copy_dir if copy_dir else output_dir
        copied_path = copy_source_file(source_path, target_dir, desired_name=original_filename)
    return GenerationResult(note_path=note_path, note_text=note_text, copied_path=copied_path)
