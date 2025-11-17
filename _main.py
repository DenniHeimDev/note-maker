#!/usr/bin/env python3
import os
import shutil
import threading
from datetime import datetime
from pathlib import Path
from typing import Optional
import tkinter as tk
from tkinter import filedialog, messagebox, scrolledtext, simpledialog, ttk

import fitz  # PyMuPDF
from dotenv import load_dotenv
from pptx import Presentation
from openai import OpenAI

AVAILABLE_MODELS = [
    "gpt-5",
    "gpt-5.1",
    "gpt-4.1",
    "gpt-4.1-mini",
    "gpt-5-mini",
]
DEFAULT_MODEL = "gpt-5.1"

LANGUAGE_OPTIONS = {
    "nynorsk": {
        "label": "Nynorsk",
        "system_prompt": (
            "Du er ein fagleg dyktig skribent som skriv klart og presist på nynorsk.\n"
            "Du får tekst frå ei fagleg presentasjon (PowerPoint eller PDF) og skal lage\n"
            "eit strukturert notat på nynorsk. Behald fagterminologi, ikkje oversett direkte til norsk vist det står på engelsk, bruk overskrifter\n"
            "og underoverskrifter der det passar, og skriv i ein stil som eignar seg\n"
            "som førebuing til undervisning eller eksamen."
        ),
        "user_template": (
            "Her er innhaldet frå presentasjonen. Lag eit strukturert notat på nynorsk\n"
            "som oppsummerer og forklarer innhaldet. Du skal ikkje referere til \"slides\"\n"
            "eller \"bilete\", berre skrive eit samanhengande notat i markdown-format.\n\n"
            "=== START AV INPUT ===\n"
            "{tekst_her}\n"
            "=== SLUTT AV INPUT ==="
        ),
        "note_suffix": "notat_nynorsk",
    },
    "bokmal": {
        "label": "Bokmål",
        "system_prompt": (
            "Du er en faglig dyktig skribent som skriver klart og presist på bokmål.\n"
            "Du får tekst fra en faglig presentasjon (PowerPoint eller PDF) og skal lage\n"
            "et strukturert notat på bokmål. Behold fagterminologi, ikke oversett direkte fra engelsk\n"
            "dersom det ikke gir mening, og bruk overskrifter og underoverskrifter der det passer."
        ),
        "user_template": (
            "Her er innholdet fra presentasjonen. Lag et strukturert notat på bokmål\n"
            "som oppsummerer og forklarer innholdet. Du skal ikke referere til \"slides\"\n"
            "eller \"bilder\", men skrive et sammenhengende notat i markdown.\n\n"
            "=== START AV INPUT ===\n"
            "{tekst_her}\n"
            "=== SLUTT AV INPUT ==="
        ),
        "note_suffix": "notat_bokmal",
    },
    "english": {
        "label": "English",
        "system_prompt": (
            "You are an expert technical writer who produces clear, structured notes in English.\n"
            "You receive text extracted from a presentation (PowerPoint or PDF) and must create\n"
            "a study note. Keep domain terminology, avoid literal translations that harm meaning,\n"
            "and use headings and subheadings where appropriate to prepare the reader for teaching or exams."
        ),
        "user_template": (
            "Here is the content from the presentation. Produce a structured note in English\n"
            "that summarizes and explains the material. Do not mention \"slides\" or \"images\";\n"
            "write a continuous markdown document instead.\n\n"
            "=== START OF INPUT ===\n"
            "{tekst_her}\n"
            "=== END OF INPUT ==="
        ),
        "note_suffix": "note_english",
    },
}
DEFAULT_LANGUAGE = "nynorsk"
LANGUAGE_LABEL_TO_KEY = {data["label"]: key for key, data in LANGUAGE_OPTIONS.items()}
LANGUAGE_LABELS = [data["label"] for data in LANGUAGE_OPTIONS.values()]

load_dotenv()


def _ensure_api_key() -> Optional[str]:
    """
    Returns the OpenAI API key from the environment, falling back to plain-text .env files.
    Supports both standard KEY=VALUE files and files that only contain the key value.
    """
    existing = os.environ.get("OPENAI_API_KEY")
    if existing:
        return existing

    candidate_paths = [
        Path(".env"),
        Path(__file__).with_name(".env"),
    ]

    for path in candidate_paths:
        if not path.exists():
            continue
        with path.open(encoding="utf-8") as handle:
            for line in handle:
                stripped = line.strip()
                if not stripped or stripped.startswith("#"):
                    continue
                if "=" in stripped:
                    _, value = stripped.split("=", 1)
                    candidate = value.strip().strip("'\"")
                else:
                    candidate = stripped
                if candidate:
                    os.environ["OPENAI_API_KEY"] = candidate
                    return candidate
    return None


OPENAI_API_KEY = _ensure_api_key()
client: Optional[OpenAI] = None
if OPENAI_API_KEY:
    client = OpenAI(api_key=OPENAI_API_KEY)


def extract_text_from_pptx(file_path: str) -> str:
    presentation = Presentation(file_path)
    sections = []
    for slide in presentation.slides:
        slide_parts = []
        for shape in slide.shapes:
            if getattr(shape, "has_text_frame", False):
                for paragraph in shape.text_frame.paragraphs:
                    text = "".join(run.text for run in paragraph.runs).strip()
                    if text:
                        slide_parts.append(text)
            elif getattr(shape, "has_table", False):
                for row in shape.table.rows:
                    cells = [cell.text.strip() for cell in row.cells if cell.text.strip()]
                    if cells:
                        slide_parts.append(" | ".join(cells))
        if slide_parts:
            sections.append("\n".join(slide_parts))
    return "\n\n".join(sections)


def extract_text_from_pdf(file_path: str) -> str:
    doc = fitz.open(file_path)
    pages = []
    try:
        for page in doc:
            page_text = page.get_text("text").strip()
            if page_text:
                pages.append(page_text)
    finally:
        doc.close()
    return "\n\n".join(pages)


def extract_text(file_path: str) -> str:
    suffix = Path(file_path).suffix.lower()
    if suffix == ".pptx":
        return extract_text_from_pptx(file_path)
    if suffix == ".pdf":
        return extract_text_from_pdf(file_path)
    raise ValueError("Ukjend filtype. Vel ei .pptx- eller .pdf-fil.")


def _language_settings(language_key: str) -> dict:
    settings = LANGUAGE_OPTIONS.get(language_key)
    if not settings:
        raise ValueError(f"Ugyldig språkval: {language_key}.")
    return settings


def generate_note_from_text(text: str, model_name: str, language_key: str) -> str:
    if not OPENAI_API_KEY:
        raise RuntimeError("Miljøvariabelen OPENAI_API_KEY er ikkje sett.")
    if not text.strip():
        raise ValueError("Fann ikkje tekst i den valde fila.")
    if client is None:
        raise RuntimeError("Fann ikkje klient for OpenAI. Kontroller API-nøkkelen.")
    if model_name not in AVAILABLE_MODELS:
        raise ValueError(f"Ugyldig modell: {model_name}.")
    settings = _language_settings(language_key)
    user_prompt = settings["user_template"].format(tekst_her=text.strip())
    response = client.chat.completions.create(
        model=model_name,
        messages=[
            {"role": "system", "content": settings["system_prompt"]},
            {"role": "user", "content": user_prompt},
        ],
    )
    note = response.choices[0].message.content
    if not note:
        raise RuntimeError("Modellen returnerte ikkje noko innhald.")
    return note.strip()


def save_note_text(note_text: str, output_dir: str, source_file: str, note_suffix: str) -> Path:
    output_path = Path(output_dir)
    output_path.mkdir(parents=True, exist_ok=True)
    file_stem = Path(source_file).stem
    note_file = output_path / f"{file_stem}_{note_suffix}.md"
    with note_file.open("w", encoding="utf-8") as handle:
        handle.write(note_text)
    return note_file


def copy_source_file(source_file: str, output_dir: str) -> Path:
    destination_dir = Path(output_dir)
    destination_dir.mkdir(parents=True, exist_ok=True)

    original_name = Path(source_file).name
    stem = Path(original_name).stem
    suffix = Path(original_name).suffix
    destination = destination_dir / original_name

    counter = 1
    while destination.exists():
        destination = destination_dir / f"{stem}_kopi_{counter}{suffix}"
        counter += 1

    shutil.copy2(source_file, destination)
    return destination


class NoteMakerApp:
    def __init__(self, root: tk.Tk) -> None:
        self.root = root
        self.root.title("Notatgenerator for forelesninger")
        self.root.geometry("900x650")

        self.input_mount = Path(os.environ.get("HOST_INPUT_DIR", "/host/input"))
        self.output_mount = Path(os.environ.get("HOST_OUTPUT_DIR", "/host/output"))
        self.copy_mount = Path(os.environ.get("HOST_COPY_DIR", "/host/copies"))
        self.file_path = tk.StringVar()
        default_output = str(
            self.output_mount if self.output_mount.exists() else Path.home()
        )
        self.output_dir = tk.StringVar(value=default_output)
        self.copy_source = tk.BooleanVar(value=False)
        default_copy_dir = self._default_copy_directory_string(default_output)
        self.copy_destination = tk.StringVar(value=default_copy_dir)
        default_model = DEFAULT_MODEL if DEFAULT_MODEL in AVAILABLE_MODELS else AVAILABLE_MODELS[0]
        self.model_name = tk.StringVar(value=default_model)
        default_language_label = LANGUAGE_OPTIONS.get(
            DEFAULT_LANGUAGE, next(iter(LANGUAGE_OPTIONS.values()))
        )["label"]
        self.language_choice = tk.StringVar(value=default_language_label)
        self.copy_source.trace_add("write", lambda *_: self._update_copy_controls_state())

        self._build_layout()
        self._update_copy_controls_state()
        self.update_api_key_status()

    def _build_layout(self) -> None:
        self.root.columnconfigure(0, weight=1)
        self.root.rowconfigure(3, weight=1)

        status_frame = ttk.LabelFrame(self.root, text="OpenAI API-nøkkel")
        status_frame.grid(row=0, column=0, padx=12, pady=(12, 6), sticky="ew")
        status_frame.columnconfigure(1, weight=1)

        ttk.Label(status_frame, text="Status:").grid(row=0, column=0, padx=6, pady=6, sticky="w")
        self.api_status_label = tk.Label(status_frame, anchor="w", font=("TkDefaultFont", 10, "bold"))
        self.api_status_label.grid(row=0, column=1, padx=6, pady=6, sticky="w")
        ttk.Label(status_frame, text="Modell:").grid(row=1, column=0, padx=6, pady=6, sticky="w")
        self.model_selector = ttk.Combobox(
            status_frame,
            textvariable=self.model_name,
            values=AVAILABLE_MODELS,
            state="readonly",
        )
        self.model_selector.grid(row=1, column=1, padx=6, pady=6, sticky="ew")
        ttk.Label(status_frame, text="Språk:").grid(row=2, column=0, padx=6, pady=6, sticky="w")
        self.language_selector = ttk.Combobox(
            status_frame,
            textvariable=self.language_choice,
            values=LANGUAGE_LABELS,
            state="readonly",
        )
        self.language_selector.grid(row=2, column=1, padx=6, pady=6, sticky="ew")

        file_frame = ttk.LabelFrame(self.root, text="1. Vel presentasjonsfil")
        file_frame.grid(row=1, column=0, padx=12, pady=6, sticky="ew")
        file_frame.columnconfigure(1, weight=1)

        ttk.Label(file_frame, text="Fil:").grid(row=0, column=0, padx=6, pady=6, sticky="w")
        self.file_entry = ttk.Entry(file_frame, textvariable=self.file_path, state="readonly")
        self.file_entry.grid(row=0, column=1, padx=6, pady=6, sticky="ew")
        ttk.Button(file_frame, text="Bla gjennom", command=self.choose_file).grid(
            row=0, column=2, padx=6, pady=6
        )

        output_frame = ttk.LabelFrame(self.root, text="2. Vel lagringsmappe")
        output_frame.grid(row=2, column=0, padx=12, pady=6, sticky="ew")
        output_frame.columnconfigure(1, weight=1)

        ttk.Label(output_frame, text="Mappe:").grid(row=0, column=0, padx=6, pady=6, sticky="w")
        self.output_entry = ttk.Entry(output_frame, textvariable=self.output_dir, state="readonly")
        self.output_entry.grid(row=0, column=1, padx=6, pady=6, sticky="ew")
        ttk.Button(output_frame, text="Vel mappe", command=self.choose_output_dir).grid(
            row=0, column=2, padx=6, pady=6
        )
        ttk.Button(
            output_frame, text="Ny mappe", command=self.create_and_select_new_folder
        ).grid(row=0, column=3, padx=6, pady=6)

        ttk.Checkbutton(
            output_frame,
            text="Kopier presentasjonen til lagringsmappa",
            variable=self.copy_source,
        ).grid(row=1, column=0, columnspan=3, padx=6, pady=(0, 6), sticky="w")

        copy_dest_frame = ttk.Frame(output_frame)
        copy_dest_frame.grid(row=2, column=0, columnspan=4, padx=6, pady=(0, 6), sticky="ew")
        copy_dest_frame.columnconfigure(1, weight=1)
        ttk.Label(copy_dest_frame, text="Kopimappe:").grid(row=0, column=0, padx=(0, 6), pady=4, sticky="w")
        self.copy_destination_entry = ttk.Entry(
            copy_dest_frame,
            textvariable=self.copy_destination,
            state="readonly",
        )
        self.copy_destination_entry.grid(row=0, column=1, padx=(0, 6), pady=4, sticky="ew")
        self.copy_dir_button = ttk.Button(
            copy_dest_frame,
            text="Vel mappe",
            command=self.choose_copy_directory,
        )
        self.copy_dir_button.grid(row=0, column=2, padx=(0, 6), pady=4)

        log_frame = ttk.LabelFrame(self.root, text="Status og logg")
        log_frame.grid(row=3, column=0, padx=12, pady=6, sticky="nsew")
        log_frame.rowconfigure(0, weight=1)
        log_frame.columnconfigure(0, weight=1)

        self.log_widget = scrolledtext.ScrolledText(
            log_frame, state="disabled", wrap="word", height=15
        )
        self.log_widget.grid(row=0, column=0, sticky="nsew", padx=6, pady=6)

        self.process_button = ttk.Button(
            self.root, text="3. Generer notat", command=self.start_processing
        )
        self.process_button.grid(row=4, column=0, padx=12, pady=(0, 6), sticky="ew")

    def update_api_key_status(self) -> None:
        has_key = bool(os.environ.get("OPENAI_API_KEY"))
        if has_key:
            symbol = "✓"
            message = "API-nøkkel funnen"
            color = "green"
        else:
            symbol = "✗"
            message = "API-nøkkel manglar"
            color = "red"
        self.api_status_label.config(text=f"{symbol} {message}", foreground=color)

    def choose_file(self) -> None:
        filetypes = [
            ("Presentasjonar (*.pptx, *.pdf)", "*.pptx *.pdf"),
            ("PowerPoint (*.pptx)", "*.pptx"),
            ("PDF (*.pdf)", "*.pdf"),
        ]
        selected = filedialog.askopenfilename(
            title="Vel ei presentasjonsfil",
            filetypes=filetypes,
            initialdir=self._initial_dir(self.input_mount),
        )
        if selected:
            self.file_path.set(selected)
            self.log_message(f"Valde fil: {selected}")

    def choose_output_dir(self) -> None:
        selected = filedialog.askdirectory(
            title="Vel lagringsmappe",
            initialdir=self._initial_dir(self.output_mount),
        )
        if selected:
            self.output_dir.set(selected)
            self.log_message(f"Valde lagringsmappe: {selected}")

    def choose_copy_directory(self) -> None:
        initial_dir = self.copy_destination.get() or self._default_copy_directory_string(self.output_dir.get())
        if initial_dir and Path(initial_dir).exists():
            start_dir = initial_dir
        else:
            start_dir = self._initial_dir(self.copy_mount)
        selected = filedialog.askdirectory(
            title="Vel kopimappe",
            initialdir=start_dir,
        )
        if selected:
            self.copy_destination.set(selected)
            self.log_message(f"Valde kopimappe: {selected}")

    def create_and_select_new_folder(self) -> None:
        parent_dir = filedialog.askdirectory(
            title="Vel foreldremappe der ny mappe skal opprettast",
            initialdir=self._initial_dir(self.output_mount),
        )
        if not parent_dir:
            return
        folder_name = simpledialog.askstring(
            "Ny mappe", "Skriv namnet på den nye mappa:", parent=self.root
        )
        if not folder_name:
            return
        new_folder = Path(parent_dir) / folder_name
        if new_folder.exists():
            messagebox.showerror(
                "Mappe finst frå før",
                f"Mappa {new_folder} finst allereie. Vel eit anna namn.",
            )
            return
        try:
            new_folder.mkdir(parents=True, exist_ok=False)
        except Exception as exc:
            messagebox.showerror("Kunne ikkje lage mappe", f"Feil ved oppretting:\n{exc}")
            return
        self.output_dir.set(str(new_folder))
        self.log_message(f"Laga og valde ny lagringsmappe: {new_folder}")

    def start_processing(self) -> None:
        file_path = self.file_path.get()
        output_dir = self.output_dir.get()
        model_name = self.model_name.get()
        language_label = self.language_choice.get()
        language_key = LANGUAGE_LABEL_TO_KEY.get(language_label)

        if not file_path:
            messagebox.showwarning("Manglar fil", "Vel ei presentasjonsfil før du held fram.")
            return
        if not Path(file_path).exists():
            messagebox.showerror("Fil finst ikkje", "Den valde fila vart ikkje funnen.")
            return
        if model_name not in AVAILABLE_MODELS:
            messagebox.showerror("Ugyldig modell", "Vel ei gyldig OpenAI-modell frå lista.")
            return
        if not language_key:
            messagebox.showerror("Ugyldig språk", "Vel eit gyldig språk frå lista.")
            return
        copy_dir = None
        copy_requested = self.copy_source.get()
        if copy_requested:
            preferred_dir = self.copy_destination.get().strip()
            copy_dir = self._copy_directory_for_session(
                output_dir, preferred_dir if preferred_dir else None
            )

        if not output_dir:
            messagebox.showwarning("Manglar mappe", "Vel lagringsmappe før du held fram.")
            return

        self.set_processing_state(True)
        self.log_message(
            f"Startar generering av notat med modell {model_name} på språk {language_label} ..."
        )
        thread = threading.Thread(
            target=self._process_workflow,
            args=(
                file_path,
                output_dir,
                copy_dir,
                copy_requested,
                model_name,
                language_key,
            ),
            daemon=True,
        )
        thread.start()

    def set_processing_state(self, is_processing: bool) -> None:
        new_state = "disabled" if is_processing else "normal"
        self.process_button.config(state=new_state)

    def _update_copy_controls_state(self) -> None:
        enabled = self.copy_source.get()
        entry_state = "readonly" if enabled else "disabled"
        button_state = "normal" if enabled else "disabled"
        self.copy_destination_entry.config(state=entry_state)
        self.copy_dir_button.config(state=button_state)

    def _process_workflow(
        self,
        file_path: str,
        output_dir: str,
        copy_dir: Optional[str],
        copy_requested: bool,
        model_name: str,
        language_key: str,
    ) -> None:
        try:
            self.log_message("Hentar tekst frå fila ...")
            extracted_text = extract_text(file_path)
            self.log_message("Tekst ekstrahert. Sender til språkmodellen ...")
            language_settings = _language_settings(language_key)
            note_text = generate_note_from_text(extracted_text, model_name, language_key)
            note_path = save_note_text(
                note_text, output_dir, file_path, language_settings["note_suffix"]
            )
            self.log_message(f"Notat lagra som: {note_path}")
            copied_path = None
            if copy_requested:
                self.log_message("Kopierer presentasjonen til eksportmappa ...")
                target_dir = copy_dir or self._copy_directory_for_session(output_dir)
                copied_path = copy_source_file(file_path, target_dir)
                self.log_message(f"Kopierte presentasjonen til {copied_path}")

            def success_message() -> None:
                message = f"Ferdig! Notatet er lagra som:\n{note_path}"
                if copied_path:
                    message += f"\nPresentasjonen vart kopiert til:\n{copied_path}"
                messagebox.showinfo("Ferdig", message)

            self.root.after(0, success_message)
        except Exception as exc:
            self.log_message(f"Feil: {exc}")
            self.root.after(
                0,
                lambda: messagebox.showerror(
                    "Feil under generering",
                    f"Arbeidsflyten stoppa på grunn av ein feil:\n{exc}",
                ),
            )
        finally:
            self.root.after(0, lambda: self.set_processing_state(False))

    def log_message(self, message: str) -> None:
        timestamp = datetime.now().strftime("%H:%M:%S")
        self.root.after(0, lambda: self._append_log(f"[{timestamp}] {message}"))

    def _initial_dir(self, mount_path: Path) -> str:
        if mount_path.exists():
            return str(mount_path)
        return os.getcwd()

    def _default_copy_directory_string(self, current_output: Optional[str] = None) -> str:
        if self.copy_mount.exists():
            return str(self.copy_mount)
        base_dir = Path(current_output) if current_output else Path(self.output_dir.get() or Path.home())
        return str(Path(base_dir) / "kopierte_presentasjonar")

    def _copy_directory_for_session(self, output_dir: str, preferred_dir: Optional[str] = None) -> str:
        if preferred_dir:
            target = Path(preferred_dir)
        elif self.copy_mount.exists():
            target = self.copy_mount
        else:
            target = Path(output_dir) / "kopierte_presentasjonar"
        target.mkdir(parents=True, exist_ok=True)
        return str(target)

    def _append_log(self, message: str) -> None:
        self.log_widget.configure(state="normal")
        self.log_widget.insert(tk.END, message + "\n")
        self.log_widget.configure(state="disabled")
        self.log_widget.see(tk.END)


def main() -> None:
    root = tk.Tk()
    NoteMakerApp(root)
    root.mainloop()


if __name__ == "__main__":
    main()
