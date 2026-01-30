"""Unit tests for note_maker.core module."""

from __future__ import annotations

from pathlib import Path

import pytest

from note_maker.core import (
    AVAILABLE_MODELS,
    DEFAULT_LANGUAGE,
    DEFAULT_MODEL,
    LANGUAGE_OPTIONS,
    extract_text,
    extract_text_from_pdf,
    extract_text_from_pptx,
    save_note_text,
)


class TestExtractTextFromPptx:
    """Tests for PowerPoint text extraction."""

    def test_extract_text_from_pptx_with_text_frames(self, tmp_path: Path) -> None:
        """Test extraction from a PPTX with text content."""
        from pptx import Presentation
        from pptx.util import Inches

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)

        left = top = Inches(1)
        width = height = Inches(2)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        tf = txBox.text_frame
        tf.text = "Hello World"

        prs.save(str(pptx_path))

        result = extract_text_from_pptx(str(pptx_path))
        assert "Hello World" in result

    def test_extract_text_from_empty_pptx(self, tmp_path: Path) -> None:
        """Test extraction from an empty PPTX."""
        from pptx import Presentation

        pptx_path = tmp_path / "empty.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        result = extract_text_from_pptx(str(pptx_path))
        assert result == ""


class TestExtractTextFromPdf:
    """Tests for PDF text extraction."""

    def test_extract_text_from_pdf_with_content(self, tmp_path: Path) -> None:
        """Test extraction from a PDF with text content."""
        import fitz

        pdf_path = tmp_path / "test.pdf"
        doc = fitz.open()
        page = doc.new_page()
        page.insert_text((100, 100), "Test PDF Content")
        doc.save(str(pdf_path))
        doc.close()

        result = extract_text_from_pdf(str(pdf_path))
        assert "Test PDF Content" in result

    def test_extract_text_from_empty_pdf(self, tmp_path: Path) -> None:
        """Test extraction from an empty PDF."""
        import fitz

        pdf_path = tmp_path / "empty.pdf"
        doc = fitz.open()
        doc.new_page()
        doc.save(str(pdf_path))
        doc.close()

        result = extract_text_from_pdf(str(pdf_path))
        assert result == ""


class TestExtractText:
    """Tests for the generic extract_text dispatcher."""

    def test_extract_text_dispatches_to_pptx(self, tmp_path: Path) -> None:
        """Test that .pptx files are routed correctly."""
        from pptx import Presentation

        pptx_path = tmp_path / "test.pptx"
        prs = Presentation()
        prs.save(str(pptx_path))

        # Should not raise
        result = extract_text(str(pptx_path))
        assert isinstance(result, str)

    def test_extract_text_dispatches_to_pdf(self, tmp_path: Path) -> None:
        """Test that .pdf files are routed correctly."""
        import fitz

        pdf_path = tmp_path / "test.pdf"
        doc = fitz.open()
        doc.new_page()
        doc.save(str(pdf_path))
        doc.close()

        # Should not raise
        result = extract_text(str(pdf_path))
        assert isinstance(result, str)

    def test_extract_text_raises_for_unsupported_format(self, tmp_path: Path) -> None:
        """Test that unsupported formats raise ValueError."""
        txt_path = tmp_path / "test.txt"
        txt_path.write_text("Some text")

        with pytest.raises(ValueError, match="Ukjend filtype"):
            extract_text(str(txt_path))


class TestSaveNoteText:
    """Tests for saving note text to files."""

    def test_save_note_text_creates_file(self, tmp_path: Path) -> None:
        """Test that note files are created correctly."""
        note_text = "# Test Note\n\nThis is a test."
        result_path = save_note_text(
            note_text,
            output_dir=tmp_path,
            source_name="presentation.pptx",
            note_suffix="notat_nynorsk",
        )

        assert result_path.exists()
        assert result_path.name == "presentation_notat_nynorsk.md"
        assert result_path.read_text(encoding="utf-8") == note_text

    def test_save_note_text_creates_output_directory(self, tmp_path: Path) -> None:
        """Test that output directory is created if it doesn't exist."""
        output_dir = tmp_path / "nested" / "output"
        note_text = "Test content"

        result_path = save_note_text(
            note_text,
            output_dir=output_dir,
            source_name="test.pdf",
            note_suffix="note_english",
        )

        assert output_dir.exists()
        assert result_path.exists()


class TestConstants:
    """Tests for module constants."""

    def test_available_models_is_not_empty(self) -> None:
        """Test that available models list is populated."""
        assert len(AVAILABLE_MODELS) > 0

    def test_default_model_is_in_available_models(self) -> None:
        """Test that default model is valid."""
        assert DEFAULT_MODEL in AVAILABLE_MODELS

    def test_language_options_contain_required_keys(self) -> None:
        """Test that language options have required structure."""
        required_keys = {"label", "system_prompt", "user_template", "note_suffix"}
        for lang_key, lang_data in LANGUAGE_OPTIONS.items():
            assert required_keys.issubset(lang_data.keys()), f"Missing keys in {lang_key}"

    def test_default_language_is_valid(self) -> None:
        """Test that default language exists in options."""
        assert DEFAULT_LANGUAGE in LANGUAGE_OPTIONS
